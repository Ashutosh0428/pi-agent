"""Data-science tools: profile a dataset and build a slide deck.

Both are **fixed, safe** tools (not arbitrary code execution): the model chooses
*what* to analyse or *what* slides to make; the tool controls *how*. That makes
them safe to expose on a public app where a raw Python runner would be an RCE.

Heavy deps (pandas / matplotlib / python-pptx) are imported lazily, so importing
this module — and running the rest of the test suite — needs none of them.
Install with ``pip install -e ".[data]"``.
"""

from __future__ import annotations

from typing import Any

from pi_agent.sandbox import Sandbox, SandboxError
from pi_agent.tools.base import Tool

MAX_ROWS = 200_000


def _resolve(path: str, sb: Sandbox):
    """Resolve a sandbox path or return an error string."""
    try:
        return sb.resolve(path), None
    except SandboxError:
        return None, f"Error: '{path}' is outside the sandbox."


def _analyze_data(args: dict[str, Any], sb: Sandbox) -> str:
    path = args.get("path")
    if not path:
        return "Error: 'path' (a CSV/Excel file in the workspace) is required."
    target, err = _resolve(path, sb)
    if err:
        return err
    if not target.is_file():
        return f"Error: '{path}' is not a file."
    try:
        import pandas as pd
    except ImportError:
        return "Error: analysis needs pandas — install with pi-agent[data]."

    try:
        if target.suffix.lower() in (".xlsx", ".xls"):
            df = pd.read_excel(target, nrows=MAX_ROWS)
        elif target.suffix.lower() in (".tsv",):
            df = pd.read_csv(target, sep="\t", nrows=MAX_ROWS)
        else:
            df = pd.read_csv(target, nrows=MAX_ROWS)
    except Exception as exc:  # noqa: BLE001 - surface any parse error to the model
        return f"Error reading '{path}': {exc}"

    lines = [f"# Data profile: {target.name}", f"Shape: {df.shape[0]} rows x {df.shape[1]} cols", ""]

    lines.append("## Columns")
    for col in df.columns:
        s = df[col]
        miss = s.isna().mean() * 100
        info = f"- **{col}** ({s.dtype}), missing {miss:.1f}%"
        if pd.api.types.is_numeric_dtype(s):
            info += f", mean={s.mean():.3g}, std={s.std():.3g}, min={s.min():.3g}, max={s.max():.3g}"
        else:
            top = s.value_counts().head(1)
            if not top.empty:
                info += f", {s.nunique()} unique, top='{top.index[0]}' ({int(top.iloc[0])})"
        lines.append(info)

    num = df.select_dtypes("number")
    if num.shape[1] >= 2:
        corr = num.corr().abs()
        pairs = []
        cols = list(corr.columns)
        for i in range(len(cols)):
            for j in range(i + 1, len(cols)):
                v = corr.iloc[i, j]
                if v == v:  # drop NaN (constant columns)
                    pairs.append((cols[i], cols[j], v))
        pairs.sort(key=lambda t: t[2], reverse=True)
        lines.append("\n## Top correlations")
        for a, b, v in pairs[:5]:
            lines.append(f"- {a} ↔ {b}: {v:.2f}")

    target_col = args.get("target")
    if target_col and target_col in df.columns:
        lines.append(f"\n## Target: {target_col}")
        ts = df[target_col]
        if pd.api.types.is_numeric_dtype(ts) and num.shape[1] >= 2:
            rel = num.corr()[target_col].drop(target_col).abs().sort_values(ascending=False)
            lines.append("Most correlated features:")
            for feat, v in rel.head(5).items():
                lines.append(f"- {feat}: {v:.2f}")
        else:
            vc = ts.value_counts(normalize=True).head(8)
            lines.append("Class balance:")
            for cls, frac in vc.items():
                lines.append(f"- {cls}: {frac:.1%}")

    return "\n".join(lines)


def _make_slides(args: dict[str, Any], sb: Sandbox) -> str:
    title = str(args.get("title") or "Presentation")
    slides = args.get("slides")
    if not isinstance(slides, list) or not slides:
        return "Error: 'slides' must be a non-empty list of {heading, bullets}."
    name = str(args.get("filename") or "deck.pptx")
    if not name.endswith(".pptx"):
        name += ".pptx"
    dest, err = _resolve(name, sb)
    if err:
        return err
    try:
        from pptx import Presentation
    except ImportError:
        return "Error: slides need python-pptx — install with pi-agent[data]."

    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title

    for sl in slides:
        if not isinstance(sl, dict):
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(sl.get("heading", ""))
        bullets = sl.get("bullets", [])
        if isinstance(bullets, str):
            bullets = [bullets]
        body = slide.placeholders[1].text_frame
        body.text = str(bullets[0]) if bullets else ""
        for b in bullets[1:]:
            body.add_paragraph().text = str(b)

    try:
        prs.save(str(dest))
    except OSError as exc:
        return f"Error saving deck: {exc}"
    return f"Created {sb.relpath(dest)} ({len(slides)} content slides). The user can download it."


def data_tools() -> list[Tool]:
    return [
        Tool(
            name="analyze_data",
            description=(
                "Profile a dataset (CSV/TSV/Excel in the workspace) like a data "
                "scientist: shape, columns, dtypes, missing %, summary stats, top "
                "correlations, and target relationship. Returns a text report to "
                "narrate insights from."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Data file path in the workspace."},
                    "target": {"type": "string", "description": "Optional target column."},
                },
                "required": ["path"],
            },
            handler=_analyze_data,
            mutating=False,
        ),
        Tool(
            name="make_slides",
            description=(
                "Build a downloadable PowerPoint (.pptx) from an outline. Provide a "
                "title and a list of slides, each with a heading and bullet points."
            ),
            input_schema={
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "filename": {"type": "string", "description": "Output .pptx name."},
                    "slides": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string"},
                                "bullets": {"type": "array", "items": {"type": "string"}},
                            },
                            "required": ["heading", "bullets"],
                        },
                    },
                },
                "required": ["title", "slides"],
            },
            handler=_make_slides,
            mutating=True,
        ),
    ]
